"""Insert 4 tech-methods slides between current slide 5 (Dataset) and slide 6 (Layer 01).

The Layer slides (6, 7, 8, ...) live on a custom slide master that python-pptx
can't cleanly clone across. So instead we add the new slides on the standard
BLANK layout (master 0) and manually paint a matching dark-slate background and
all shape chrome — the visual style still matches the existing layer slides.
Finally, reorder the slide list so the 4 new slides sit at positions 6..9, and
renumber the page-counter footer on every slide (now N / 15).
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt, Inches
from pptx.oxml.ns import qn
from lxml import etree

SRC = Path("Massive Data Final Project.pptx")
DST = Path("Massive Data Final Project.pptx")  # overwrite in place

# ----- color palette (reuse from existing layer slides) -----
ACCENT_TECH = RGBColor(0x0E, 0xA5, 0xE9)     # cyan-500 — tech methods accent
LABEL_BLUE  = RGBColor(0x1C, 0x72, 0x93)     # existing top-label color
TITLE_DARK  = RGBColor(0x1E, 0x29, 0x3B)
SUBTITLE    = RGBColor(0x47, 0x55, 0x69)
META_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
SOFT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
GREEN_OK    = RGBColor(0x05, 0x96, 0x69)
YELLOW_HIL  = RGBColor(0xF4, 0xB9, 0x42)
CARD_FILL   = RGBColor(0x1E, 0x29, 0x3B)     # dark slate panel — matches the layer cards
CARD_FILL_2 = RGBColor(0x0F, 0x17, 0x2A)     # darker variant
SLIDE_BG    = RGBColor(0xFF, 0xFF, 0xFF)     # white — matches original layer slides
LIGHT_CARD  = RGBColor(0xF1, 0xF5, 0xF9)     # very light gray panel for body cards on white


# ---------------- helpers ----------------

def emu(inches: float) -> int:
    return int(inches * 914400)


def add_blank_slide(prs):
    """Add a slide on the layer-slide master/layout (master 2's DEFAULT layout)
    so the dark theme + lack of Georgetown watermark match the existing layer slides.
    Then strip placeholder shapes and overlay an explicit dark background fill."""
    layer_layout = prs.slide_masters[2].slide_layouts[0]  # DEFAULT (slideLayout15)
    new_slide = prs.slides.add_slide(layer_layout)

    # Remove any default placeholder shapes
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)

    # Set the slide background to a solid dark fill (matches layer cards palette).
    cSld = new_slide._element.cSld
    existing_bg = cSld.find(qn("p:bg"))
    if existing_bg is not None:
        cSld.remove(existing_bg)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="{0:02X}{1:02X}{2:02X}"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>'
    ).format(SLIDE_BG[0], SLIDE_BG[1], SLIDE_BG[2])
    bg_el = etree.fromstring(bg_xml)
    cSld.insert(0, bg_el)

    return new_slide


def move_slide(prs, slide, new_idx: int):
    """Reorder a slide to position `new_idx` (0-based) in sldIdLst."""
    # Resolve the rId by matching target_part against the slide's part.
    rId = None
    for rel in prs.part.rels.values():
        if rel.target_part is slide.part:
            rId = rel.rId
            break
    if rId is None:
        raise RuntimeError("could not locate rId for slide")

    sldIdLst = prs.slides._sldIdLst
    target = None
    for s in list(sldIdLst):
        if s.get(qn("r:id")) == rId:
            target = s
            break
    if target is None:
        raise RuntimeError("sldId entry not found")

    sldIdLst.remove(target)
    sldIdLst.insert(new_idx, target)


def clear_shapes(slide):
    spTree = slide.shapes._spTree
    # keep nvGrpSpPr and grpSpPr (the first two children); drop the rest
    keep_tags = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    for child in list(spTree):
        if child.tag not in keep_tags:
            spTree.remove(child)


def add_text_box(slide, x_in, y_in, w_in, h_in, *, fill=None, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(x_in), emu(y_in), emu(w_in), emu(h_in)
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top  = Pt(2); tf.margin_bottom = Pt(2)
    return shp


def set_text(shape, runs, *, align=None, anchor=None):
    """runs: list of (text, dict(font_name, size_pt, bold, color, italic)) — strings get '\n' as new paragraph."""
    tf = shape.text_frame
    # purge default paragraph except the first
    while len(tf.paragraphs) > 1:
        tf._txBody.remove(tf.paragraphs[-1]._p)
    p = tf.paragraphs[0]
    # clear existing runs
    for r in list(p.runs):
        p._p.remove(r._r)
    if align is not None:
        p.alignment = align
    if anchor is not None:
        tf.vertical_anchor = anchor
    paras = [p]
    cur = p
    for entry in runs:
        if entry == "\n":
            new_p = tf.add_paragraph()
            if align is not None:
                new_p.alignment = align
            cur = new_p
            continue
        text, fmt = entry
        run = cur.add_run()
        run.text = text
        f = run.font
        if fmt.get("font"):  f.name = fmt["font"]
        if fmt.get("size"):  f.size = Pt(fmt["size"])
        if fmt.get("bold") is not None: f.bold = fmt["bold"]
        if fmt.get("italic") is not None: f.italic = fmt["italic"]
        if fmt.get("color"): f.color.rgb = fmt["color"]


def add_line(slide, x1_in, y1_in, x2_in, y2_in, color=META_GRAY, width_pt=0.75):
    line = slide.shapes.add_connector(1, emu(x1_in), emu(y1_in), emu(x2_in), emu(y2_in))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


# ---------------- common page chrome ----------------

def add_chrome(slide, top_label, title, subtitle, footer_meta, page_str):
    # top label (left aligned, top of slide)
    s = add_text_box(slide, 0.50, 0.30, 9.00, 0.30)
    set_text(s, [(top_label, dict(font="Calibri", size=10, bold=True, color=LABEL_BLUE))])
    # title — dark navy (matches layer slides)
    s = add_text_box(slide, 0.50, 0.62, 9.20, 0.75)
    set_text(s, [(title, dict(font="Georgia", size=28, bold=True, color=TITLE_DARK))])
    # subtitle — slate gray
    s = add_text_box(slide, 0.50, 1.40, 9.20, 0.35)
    set_text(s, [(subtitle, dict(font="Calibri", size=13, color=SUBTITLE))])
    # footer meta
    s = add_text_box(slide, 0.50, 5.20, 9.00, 0.25)
    set_text(s, [(footer_meta, dict(font="Calibri", size=9, color=META_GRAY))])
    # page number
    s = add_text_box(slide, 8.90, 5.40, 0.70, 0.20)
    set_text(s, [(page_str, dict(font="Calibri", size=8, color=META_GRAY))])


# ---------------- slide builders ----------------

def build_pipeline_slide(slide):
    add_chrome(
        slide,
        "METHODS  —  PIPELINE ARCHITECTURE",
        "End-to-End Reproducible Pipeline",
        "One input → one output per stage  ·  Parquet at every boundary  ·  any stage re-runnable",
        "scripts/run_pipeline.py orchestrates 5 stages  ·  61,452 employer rows  →  3,968 merged cities",
        "6 / 15",
    )

    # LEFT: data flow boxes (vertical chain)
    base_x = 0.50
    base_y = 1.95
    box_w  = 4.40
    box_h  = 0.42
    gap    = 0.08

    stages = [
        ("RAW",       "data/raw/h1b_fy2024.csv",            "USCIS FY2024 employer release · UTF-16 · 61,452 rows", LABEL_BLUE),
        ("CLEAN",     "scripts/cleaning.py",                "snake_case · NAICS split · zip3 proxy · re-derived approval_rate", ACCENT_TECH),
        ("FEATURES",  "scripts/feature_engineering.py",     "is_stem · region · size_bucket · share_continuing · log volumes", ACCENT_TECH),
        ("MERGE",     "scripts/merge_external.py  (+ _spark.py)", "USCIS × LCA × Census  →  3,968 cities  ·  96.7% petition coverage", ACCENT_TECH),
        ("MODEL",     "modeling.py  +  modeling_city.py",   "OLS / RandomForest / XGBoost  ·  structural vs enriched specs", ACCENT_TECH),
        ("ANALYSIS",  "scripts/analysis.py",                "8 figures + variance decomposition tables → outputs/", ACCENT_TECH),
    ]

    y = base_y
    for tag, script, descr, color in stages:
        # tag pill
        pill = add_text_box(slide, base_x, y, 0.85, box_h, fill=CARD_FILL_2)
        set_text(pill, [(tag, dict(font="Calibri", size=9, bold=True, color=color))],
                 align=2)  # center
        # script box
        sbox = add_text_box(slide, base_x + 0.90, y, box_w - 0.90, box_h, fill=CARD_FILL)
        set_text(sbox, [
            (script, dict(font="Consolas", size=10, bold=True, color=WHITE)),
            ("\n", {}),
            (descr, dict(font="Calibri", size=8.5, color=SOFT_GRAY)),
        ])
        y += box_h + gap

    # arrows between boxes (small triangle markers at right side)
    # Use connectors centered to the boxes
    for i in range(len(stages) - 1):
        ax = base_x + 0.42
        ay1 = base_y + (i + 1) * (box_h + gap) - gap + 0.01
        ay2 = ay1 + gap - 0.02
        add_line(slide, ax, ay1, ax, ay2, color=ACCENT_TECH, width_pt=1.25)

    # RIGHT: callout cards
    rx = 5.10
    ry = 1.95
    rw = 4.40

    # Card 1 — storage
    c1 = add_text_box(slide, rx, ry, rw, 1.05, fill=CARD_FILL)
    set_text(c1, [], )
    lab = add_text_box(slide, rx + 0.20, ry + 0.10, rw - 0.40, 0.25)
    set_text(lab, [("STORAGE  ·  PARQUET EVERYWHERE", dict(font="Calibri", size=10, bold=True, color=ACCENT_TECH))])
    body = add_text_box(slide, rx + 0.20, ry + 0.35, rw - 0.40, 0.65)
    set_text(body, [
        ("Columnar · compressed · schema-typed.", dict(font="Calibri", size=11, color=WHITE)),
        ("\n", {}),
        ("A full decade of releases (~1M rows) still fits on a laptop.",
         dict(font="Calibri", size=10, color=SOFT_GRAY)),
    ])

    # Card 2 — modularity
    ry2 = ry + 1.15
    c2 = add_text_box(slide, rx, ry2, rw, 1.05, fill=CARD_FILL)
    lab = add_text_box(slide, rx + 0.20, ry2 + 0.10, rw - 0.40, 0.25)
    set_text(lab, [("MODULARITY  ·  1-IN  /  1-OUT", dict(font="Calibri", size=10, bold=True, color=ACCENT_TECH))])
    body = add_text_box(slide, rx + 0.20, ry2 + 0.35, rw - 0.40, 0.65)
    set_text(body, [
        ("Each stage reads one parquet, writes one parquet.",
         dict(font="Calibri", size=11, color=WHITE)),
        ("\n", {}),
        ("Re-run any stage independently after editing — no full rebuild.",
         dict(font="Calibri", size=10, color=SOFT_GRAY)),
    ])

    # Card 3 — orchestration
    ry3 = ry2 + 1.15
    c3 = add_text_box(slide, rx, ry3, rw, 1.00, fill=CARD_FILL)
    lab = add_text_box(slide, rx + 0.20, ry3 + 0.10, rw - 0.40, 0.25)
    set_text(lab, [("ORCHESTRATION  ·  THIN SHELL", dict(font="Calibri", size=10, bold=True, color=ACCENT_TECH))])
    body = add_text_box(slide, rx + 0.20, ry3 + 0.35, rw - 0.40, 0.55)
    set_text(body, [
        ("python scripts/run_pipeline.py",
         dict(font="Consolas", size=11, bold=True, color=YELLOW_HIL)),
        ("\n", {}),
        ("Swap to Airflow / Prefect / Dagster — only the trigger changes, not the stages.",
         dict(font="Calibri", size=10, color=SOFT_GRAY)),
    ])


def build_cleaning_slide(slide):
    add_chrome(
        slide,
        "METHODS  —  CLEANING & FEATURE ENGINEERING",
        "Honest Features from a Messy Public Release",
        "Feature set deliberately bounded to what USCIS actually exposes",
        "61,452 employer rows · 17 raw columns → 25 modeling columns · zero leakage from target",
        "7 / 15",
    )

    # LEFT card — Cleaning decisions
    lx, ly, lw, lh = 0.50, 1.95, 4.40, 3.10
    panel = add_text_box(slide, lx, ly, lw, lh, fill=CARD_FILL)
    lab = add_text_box(slide, lx + 0.20, ly + 0.10, lw - 0.40, 0.30)
    set_text(lab, [("CLEANING DECISIONS", dict(font="Calibri", size=10, bold=True, color=ACCENT_TECH))])

    rows_left = [
        ("UTF-16 LE w/ BOM",       "encoding fix in cleaning.py"),
        ("NAICS string split",     "\"54 - Professional...\" → code + label"),
        ("Zip → 5-digit + zip3",   "zip3 = region proxy (Layer 2)"),
        ("Re-derive approval_rate","raw nudges 0/1 by ~10⁻⁶ — distorts means"),
        ("LCA wage cap $500k/yr",  "raw values reach $238M (hourly/annual mix)"),
        ("Keep \"UNKNOWN\" employer", "kept as distinct category, not dropped"),
    ]

    row_h = 0.42
    yy = ly + 0.45
    for k, v in rows_left:
        kbox = add_text_box(slide, lx + 0.20, yy, 1.85, row_h)
        set_text(kbox, [(k, dict(font="Calibri", size=10.5, bold=True, color=WHITE))],
                 anchor=3)  # middle
        vbox = add_text_box(slide, lx + 2.05, yy, lw - 2.25, row_h)
        set_text(vbox, [(v, dict(font="Calibri", size=9.5, color=SOFT_GRAY))], anchor=3)
        # divider line
        add_line(slide, lx + 0.20, yy + row_h, lx + lw - 0.20, yy + row_h,
                 color=RGBColor(0x33, 0x41, 0x55), width_pt=0.5)
        yy += row_h

    # RIGHT card — Engineered features
    rx, ry, rw, rh = 5.10, 1.95, 4.40, 3.10
    panel = add_text_box(slide, rx, ry, rw, rh, fill=CARD_FILL)
    lab = add_text_box(slide, rx + 0.20, ry + 0.10, rw - 0.40, 0.30)
    set_text(lab, [("ENGINEERED FEATURES", dict(font="Calibri", size=10, bold=True, color=ACCENT_TECH))])

    # Structural sub-section
    sub1 = add_text_box(slide, rx + 0.20, ry + 0.42, rw - 0.40, 0.25)
    set_text(sub1, [("STRUCTURAL  ·  USCIS only", dict(font="Calibri", size=9, bold=True, color=GREEN_OK))])

    structural = [
        ("is_stem",            "NAICS prefix ∈ {51, 54, 33, 52, 61}"),
        ("region",             "Census 4-region map"),
        ("size_bucket",        "single / small / medium / large / mega"),
        ("share_continuing",   "application mix"),
        ("log_petitions, log_employer_total", "volume controls"),
    ]
    yy = ry + 0.70
    for k, v in structural:
        kbox = add_text_box(slide, rx + 0.20, yy, 2.30, 0.22)
        set_text(kbox, [(k, dict(font="Consolas", size=9, bold=True, color=WHITE))], anchor=3)
        vbox = add_text_box(slide, rx + 2.55, yy, rw - 2.75, 0.22)
        set_text(vbox, [(v, dict(font="Calibri", size=9, color=SOFT_GRAY))], anchor=3)
        yy += 0.24

    # divider
    add_line(slide, rx + 0.20, yy + 0.05, rx + rw - 0.20, yy + 0.05,
             color=RGBColor(0x33, 0x41, 0x55), width_pt=0.5)
    yy += 0.12

    # Enrichment sub-section
    sub2 = add_text_box(slide, rx + 0.20, yy, rw - 0.40, 0.22)
    set_text(sub2, [("ENRICHMENT  ·  LCA + Census  ·  city-level only",
                     dict(font="Calibri", size=9, bold=True, color=YELLOW_HIL))])
    yy += 0.26

    enrichment = [
        ("log_wage",                  "OFLC LCA avg annual wage"),
        ("log_population, log_density","Census pop + density"),
        ("log_rent, employed_pct",    "BLS housing + labor"),
        ("unemployed_z, log_lca",     "labor slack + LCA volume"),
    ]
    for k, v in enrichment:
        kbox = add_text_box(slide, rx + 0.20, yy, 2.30, 0.20)
        set_text(kbox, [(k, dict(font="Consolas", size=9, bold=True, color=WHITE))], anchor=3)
        vbox = add_text_box(slide, rx + 2.55, yy, rw - 2.75, 0.20)
        set_text(vbox, [(v, dict(font="Calibri", size=9, color=SOFT_GRAY))], anchor=3)
        yy += 0.22


def build_spark_slide(slide):
    add_chrome(
        slide,
        "SCALABILITY  —  DISTRIBUTED COMPUTE",
        "Same Merge, Pandas and PySpark",
        "merge_external_spark.py: cluster-ready by design, verified bit-identical on local[*]",
        "PySpark 3.5  ·  Spark SQL adaptive execution  ·  one-line swap from local[*] to Dataproc / EMR / Databricks",
        "8 / 15",
    )

    # LEFT: code skeleton card
    lx, ly, lw, lh = 0.50, 1.95, 4.90, 3.10
    panel = add_text_box(slide, lx, ly, lw, lh, fill=CARD_FILL_2)
    lab = add_text_box(slide, lx + 0.20, ly + 0.10, lw - 0.40, 0.25)
    set_text(lab, [("merge_external_spark.py  ·  EXCERPT",
                    dict(font="Calibri", size=10, bold=True, color=ACCENT_TECH))])

    code = (
        "spark = (SparkSession.builder\n"
        "    .appName(\"h1b-merge-external\")\n"
        "    .config(\"spark.sql.adaptive.enabled\", \"true\")\n"
        "    .config(\"spark.sql.shuffle.partitions\", \"16\")\n"
        "    .getOrCreate())\n"
        "\n"
        "uscis_by_city = (df.groupBy(\"state\",\"city\")\n"
        "  .agg(F.sum(\"total_petitions\"),\n"
        "       F.countDistinct(\"employer\"),\n"
        "       F.avg(\"is_stem\"),\n"
        "       F.first(\"region\")))\n"
        "\n"
        "merged = uscis.join(external,\n"
        "    on=[\"state\",\"city\"], how=\"inner\")\n"
        "\n"
        "# cluster: merged.write.parquet(OUT)\n"
        "# local : merged.toPandas().to_parquet(OUT)"
    )
    code_box = add_text_box(slide, lx + 0.20, ly + 0.40, lw - 0.40, lh - 0.55, fill=BLACK)
    set_text(code_box, [(code, dict(font="Consolas", size=9, color=RGBColor(0x9C, 0xDC, 0xFE)))])

    # RIGHT TOP: parity verification
    rx, ry, rw = 5.60, 1.95, 3.90
    panel = add_text_box(slide, rx, ry, rw, 1.55, fill=CARD_FILL)
    lab = add_text_box(slide, rx + 0.20, ry + 0.10, rw - 0.40, 0.25)
    set_text(lab, [("PARITY VERIFICATION  ·  PANDAS  vs  SPARK",
                    dict(font="Calibri", size=10, bold=True, color=GREEN_OK))])

    rows = [
        ("Output rows",        "3,968  ≡  3,968"),
        ("approval_rate Δ",    "0.0"),
        ("avg_wage Δ",         "≤ 8e-13  (float roundoff)"),
    ]
    yy = ry + 0.45
    for k, v in rows:
        kb = add_text_box(slide, rx + 0.20, yy, 1.80, 0.30)
        set_text(kb, [(k, dict(font="Calibri", size=10.5, color=WHITE))], anchor=3)
        vb = add_text_box(slide, rx + 2.00, yy, rw - 2.20, 0.30)
        set_text(vb, [(v, dict(font="Consolas", size=11, bold=True, color=GREEN_OK))], anchor=3)
        add_line(slide, rx + 0.20, yy + 0.30, rx + rw - 0.20, yy + 0.30,
                 color=RGBColor(0x33, 0x41, 0x55), width_pt=0.5)
        yy += 0.32

    # RIGHT BOTTOM: deploy story
    ry2 = ry + 1.65
    panel = add_text_box(slide, rx, ry2, rw, 1.45, fill=CARD_FILL)
    lab = add_text_box(slide, rx + 0.20, ry2 + 0.10, rw - 0.40, 0.25)
    set_text(lab, [("ONE-LINE SCALE-OUT",
                    dict(font="Calibri", size=10, bold=True, color=YELLOW_HIL))])

    body = add_text_box(slide, rx + 0.20, ry2 + 0.40, rw - 0.40, 0.55)
    set_text(body, [
        (".toPandas().to_parquet(OUT)", dict(font="Consolas", size=10, bold=True, color=SOFT_GRAY)),
        ("\n", {}),
        ("→  ", dict(font="Calibri", size=11, color=YELLOW_HIL)),
        (".write.parquet(OUT)", dict(font="Consolas", size=10, bold=True, color=WHITE)),
    ])
    body2 = add_text_box(slide, rx + 0.20, ry2 + 0.95, rw - 0.40, 0.45)
    set_text(body2, [
        ("Same DataFrame code targets Dataproc / EMR / Databricks unchanged.",
         dict(font="Calibri", size=9.5, color=SOFT_GRAY)),
    ])


def build_modeling_slide(slide):
    add_chrome(
        slide,
        "METHODS  —  MODELING SETUP",
        "Three Model Families, One Ceiling Question",
        "Goal is to bound explanatory power, not to search hyperparameters",
        "ColumnTransformer + sparse one-hot (min_frequency = 20)  ·  80/20 split  ·  random_state = 42  ·  R² + MAE on test",
        "9 / 15",
    )

    # LEFT card — preprocessing pipeline
    lx, ly, lw, lh = 0.50, 1.95, 4.40, 3.10
    panel = add_text_box(slide, lx, ly, lw, lh, fill=CARD_FILL)
    lab = add_text_box(slide, lx + 0.20, ly + 0.10, lw - 0.40, 0.25)
    set_text(lab, [("PREPROCESSING PIPELINE",
                    dict(font="Calibri", size=10, bold=True, color=ACCENT_TECH))])

    # vertical chain inside the card
    steps = [
        ("X  =  numeric  +  categorical",  "share_continuing, log_*  +  state, naics, region, size_bucket"),
        ("ColumnTransformer",              "num: passthrough  ·  cat: OneHotEncoder(min_frequency=20, sparse)"),
        ("sklearn Pipeline",               "fit on train  ·  predict on held-out"),
        ("train_test_split (80 / 20)",     "random_state = 42  ·  identical rows across model families"),
        ("Score on test",                  "R²  +  MAE  ·  reported as the ceiling, not a tuning target"),
    ]
    yy = ly + 0.45
    box_h = 0.46
    for i, (head, sub) in enumerate(steps):
        b = add_text_box(slide, lx + 0.25, yy, lw - 0.50, box_h, fill=CARD_FILL_2)
        h = add_text_box(slide, lx + 0.40, yy + 0.04, lw - 0.80, 0.22)
        set_text(h, [(head, dict(font="Calibri", size=10.5, bold=True, color=WHITE))])
        s = add_text_box(slide, lx + 0.40, yy + 0.24, lw - 0.80, 0.22)
        set_text(s, [(sub, dict(font="Calibri", size=8.5, color=SOFT_GRAY))])
        yy += box_h + 0.02

    # RIGHT card top — three model families
    rx, ry, rw = 5.10, 1.95, 4.40
    rh = 1.85
    panel = add_text_box(slide, rx, ry, rw, rh, fill=CARD_FILL)
    lab = add_text_box(slide, rx + 0.20, ry + 0.10, rw - 0.40, 0.25)
    set_text(lab, [("THREE MODEL FAMILIES",
                    dict(font="Calibri", size=10, bold=True, color=ACCENT_TECH))])

    models = [
        ("OLS",            "LinearRegression  —  canonical baseline"),
        ("Random Forest",  "n=200–300, max_depth=12, n_jobs=-1, seed=42"),
        ("XGBoost",        "n=400–500, depth=5–6, lr=0.05, tree_method=hist"),
    ]
    yy = ry + 0.45
    for name, cfg in models:
        nb = add_text_box(slide, rx + 0.20, yy, 1.50, 0.40)
        set_text(nb, [(name, dict(font="Calibri", size=12, bold=True, color=WHITE))], anchor=3)
        cb = add_text_box(slide, rx + 1.70, yy, rw - 1.90, 0.40)
        set_text(cb, [(cfg, dict(font="Consolas", size=9, color=SOFT_GRAY))], anchor=3)
        add_line(slide, rx + 0.20, yy + 0.40, rx + rw - 0.20, yy + 0.40,
                 color=RGBColor(0x33, 0x41, 0x55), width_pt=0.5)
        yy += 0.42

    # RIGHT card bottom — two specs
    ry2 = ry + rh + 0.10
    rh2 = 3.10 - rh - 0.10
    panel = add_text_box(slide, rx, ry2, rw, rh2, fill=CARD_FILL_2)
    lab = add_text_box(slide, rx + 0.20, ry2 + 0.08, rw - 0.40, 0.25)
    set_text(lab, [("TWO SPECS  ·  IDENTICAL ROWS  →  FAIR Δ R²",
                    dict(font="Calibri", size=10, bold=True, color=YELLOW_HIL))])

    body = add_text_box(slide, rx + 0.20, ry2 + 0.34, rw - 0.40, rh2 - 0.40)
    set_text(body, [
        ("Structural", dict(font="Calibri", size=10, bold=True, color=GREEN_OK)),
        (":  state · region · share_stem · share_continuing · log_petitions · log_employers",
         dict(font="Calibri", size=9.5, color=SOFT_GRAY)),
        ("\n", {}),
        ("Enriched",  dict(font="Calibri", size=10, bold=True, color=YELLOW_HIL)),
        (":  + log_wage · log_population · log_density · log_rent · employed_pct · unemployed_z · log_lca",
         dict(font="Calibri", size=9.5, color=SOFT_GRAY)),
    ])


# ---------------- main ----------------

def main():
    prs = Presentation(str(SRC))

    # 1) Add 4 fresh blank slides with dark background.
    new_slides = [add_blank_slide(prs) for _ in range(4)]

    # 2) Populate them
    build_pipeline_slide(new_slides[0])
    build_cleaning_slide(new_slides[1])
    build_spark_slide(new_slides[2])
    build_modeling_slide(new_slides[3])

    # 3) Move them in sldIdLst to positions 5, 6, 7, 8 (0-based) so they sit
    #    after slide 5 (Dataset, idx 4) and before original Layer 01 (idx 5).
    target_positions = [5, 6, 7, 8]
    for slide, pos in zip(new_slides, target_positions):
        move_slide(prs, slide, pos)

    # 4) Renumber the page footers on the *presentation* slides only (1..15).
    #    The original deck had 11 presentation slides + 13 unused template
    #    slides at the tail; the original footer read "X / 11". We now have
    #    15 presentation slides; template tail (16..28) keeps no footer.
    import re
    PRESENTATION_LEN = 15
    pat = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")
    for i, slide in enumerate(prs.slides, start=1):
        if i > PRESENTATION_LEN:
            break
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            for para in shp.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs)
                if pat.match(txt):
                    runs = list(para.runs)
                    runs[0].text = f"{i} / {PRESENTATION_LEN}"
                    for r in runs[1:]:
                        r.text = ""

    prs.save(str(DST))
    print(f"saved {DST}  ·  total slides = {len(prs.slides)}")


if __name__ == "__main__":
    main()
